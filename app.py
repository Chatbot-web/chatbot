from flask import Flask, request, jsonify
import boto3
from botocore.exceptions import ClientError
from flask_cors import CORS
import json

import tempfile
import os
import fitz  # PyMuPDF
import docx
import openpyxl
from pptx import Presentation
import base64
import os 

# Initialize Flask app
app = Flask(__name__)
CORS(app)  # Enable CORS for all routes


bedrock_runtime = boto3.client(
    "bedrock-runtime",
    aws_access_key_id=os.environ.get("AWS_ACCESS_KEY"),
    aws_secret_access_key=os.environ.get("AWS_SECRET_KEY"),
    region_name=os.environ.get("REGION_NAME")
)

@app.route('/get_code', methods=['POST'])
def get_code():
    try:
        # Parse conversation history and file from POST request
        conversation_history = json.loads(request.form.get('user_message'))
        file = request.files.get('file')

        if not conversation_history:
            return jsonify({"error": "Invalid input, 'user_message' is required."}), 400

        content = ""
        image_base_64 = ""
        if file:
            filename = file.filename
            temp_dir = tempfile.gettempdir()  # Get the system's temporary directory
            file_path = os.path.join(temp_dir, filename)
            
            # Ensure the directory exists
            os.makedirs(os.path.dirname(file_path), exist_ok=True)
            
            file.save(file_path)

            _, file_extension = os.path.splitext(filename)
            print(file_extension, "extension")

            if(file_extension.lower() in ['.jpeg', '.jpg', '.png']):
                print("in image")
                image_base_64 = extract_base64_from_image(file_path)

            else:
                content = extract_content(file_path, filename)
                print("content", content)
                
            os.remove(file_path)  # Clean up the temporary file

        # Prepare the prompt
        prompt = """You are a versatile AI assistant in a chatbot interface. Your role is to:
        1. Analyze the entire conversation history provided to maintain context.
        2. Pay special attention to the most recent user query, but consider previous interactions for relevance.
        3. When referring to past conversations, be explicit about which part you're referencing.
        4. Maintain a professional and helpful tone throughout the conversation.
        5. If you're unsure about any information, state that clearly rather than making assumptions.
        6. Provide concise yet comprehensive answers, breaking down complex information when necessary.
        7. Answer a wide variety of questions on any topic.
        8. Provide code snippets in any programming language when requested, ensuring proper formatting.
        9. Analyze the full conversation history for context when needed.
        10. Respond concisely yet comprehensively to the most recent query.
        11. Ask for clarification if a question is unclear.
        12. Maintain a helpful and professional tone.

        Note: Always strive for accuracy and clarity in your responses.

        Respond to the user's most recent question, using context from previous exchanges if relevant."""


        # For image passing scenarios
        if(image_base_64 != ""):
            print("in image prompt")
            messages = [
                    {
                        "role": "user",
                        "content": [
                            {"image": {"format": "png", "source": {"bytes": image_base_64}}},
                            {"text": f"{conversation_history} - {prompt}"}
                        ],
                    }
                ]
            print("message", messages[0]["content"][1]["text"])
        else:
            print("in normal prompt")
            # Prepare the messages for the Bedrock model
            messages = [
                {
                    "role": "user",
                    "content": [
                        {"text": f"{conversation_history} - {prompt}"}
                    ],
                }
            ]
            print("message", messages[0]["content"][0]["text"])
        if(content != ""):
            updated_msg = messages[0]["content"][0]["text"]
            messages[0]["content"][0]["text"] = f"{updated_msg} - File content: {content}"
            print("message", messages[0]["content"][0]["text"])

        # Interact with the Bedrock model
        response = bedrock_runtime.converse(
            modelId=os.environ.get("MODEL_ID"),
            messages=messages
        )

        # Extract and return the response content
        response_text = response["output"]["message"]["content"][0]["text"]
        print("---------------------------------------------------------------------------------------")
        print(response_text)
        return jsonify({"response": response_text}), 200

    except ClientError as e:
        print({"error": f"ClientError: {e.response['Error']['Message']}"})
        return jsonify({"error": f"ClientError: {e.response['Error']['Message']}"}), 500
    except Exception as e:
        print({"error": str(e)})
        return jsonify({"error": str(e)}), 500

def extract_content(file_path, filename):
    _, file_extension = os.path.splitext(filename)
    
    if file_extension.lower() in ['.pdf']:
        return extract_pdf(file_path)
    elif file_extension.lower() in ['.docx', '.doc']:
        return extract_word(file_path)
    elif file_extension.lower() in ['.xlsx', '.xls']:
        return extract_excel(file_path)
    elif file_extension.lower() in ['.pptx', '.ppt']:
        return extract_powerpoint(file_path)
    elif file_extension.lower() == '.txt':
        return extract_txt(file_path)
    elif file_extension.lower() == '.json':
        return extract_json(file_path)
    else:
        return "Unsupported file type"

def extract_pdf(file_path):
    content = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            content += page.get_text()
    return content

def extract_word(file_path):
    doc = docx.Document(file_path)
    content = "\n".join([para.text for para in doc.paragraphs])
    return content

def extract_excel(file_path):
    workbook = openpyxl.load_workbook(file_path)
    content = ""
    for sheet in workbook:
        for row in sheet.iter_rows(values_only=True):
            content += " ".join(str(cell) for cell in row if cell is not None) + "\n"
    return content

def extract_powerpoint(file_path):
    prs = Presentation(file_path)
    content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                content += shape.text + "\n"
    return content

def extract_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    return content

def extract_json(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        data = json.load(file)
    return json.dumps(data, indent=2)  # Returns a formatted JSON string

def extract_base64_from_image(file_path):
    print("in image function")
    with open(file_path, "rb") as f:
        content = f.read()
    return content

if __name__ == '__main__':
    app.run(debug=True)
